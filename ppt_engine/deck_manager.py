# ppt_engine/deck_manager.py

import time
import copy
from pptx import Presentation
from pptx.slide import Slide

def open_ppt(file_path: str):
    """
    打开并返回 Presentation 对象 (python-pptx).
    相当于原先 win32com 的 prs = ppt_app.Presentations.Open(...)
    """
    prs = Presentation(file_path)
    return prs

def save_ppt(prs, output_path: str):
    """
    保存到指定路径 (相当于原先 prs.SaveAs(output_path, 24))
    """
    prs.save(output_path)

def close_ppt(prs):
    """
    python-pptx没有专门的关闭函数；可留空或记录日志
    """
    pass

def copy_slide_after(prs, base_index: int, count: int = 1, sleep_time=0.5):
    """
    在 base_index 这张幻灯片后复制 count 次 (python-pptx实现).
    返回复制后幻灯片总数.

    说明:
    - python-pptx无内置复制slide的API，这里手写一个简易clone逻辑:
      1) 找到 slides[base_index-1] (因为 python-pptx是0-based)
      2) 复制其 shapes/notes(可选)，然后插入在 base_index 位置后面
      3) 重复 count 次
    - 可能无法完整复制动画/音视频/SmartArt等高级要素
    - sleep_time 仅用来模拟在原先 win32com 逻辑里给粘贴留时间(在纯 python-pptx 中不一定需要)
    """
    # python-pptx: slides 是 0-based，而你的映射或逻辑多半是 1-based
    # 这里假设你传进来的 base_index 也是 1-based
    # => 转成 python-pptx 的下标:
    actual_idx = base_index - 1

    # 如果超出范围，直接返回
    if actual_idx < 0 or actual_idx >= len(prs.slides):
        return len(prs.slides)

    slide_to_clone = prs.slides[actual_idx]

    for _ in range(count):
        time.sleep(sleep_time)  # 模拟原先 win32com 中的复制粘贴等待
        clone_slide(prs, slide_to_clone)

    return len(prs.slides)

def clone_slide(prs: Presentation, src_slide: Slide):
    """
    在 python-pptx 里克隆 src_slide，并将新幻灯片
    紧跟在 src_slide 后面。

    步骤:
      1) 新增一张空白幻灯片(默认使用最后一个布局,例如空白布局).
      2) 用XML方式复制 src_slide 的 spTree(所有 shapes) 到新幻灯片.
      3) 将新幻灯片移动到 src_slide 后面(下标 = src_idx+1).

    注意:
      - 不支持复制动画、SmartArt、OLE 等复杂要素.
      - 也未复制notes, chart, media. 如需可自行补充.
    """
    # 1) 新增空白幻灯片(默认选用最后一个layout)
    blank_layout = prs.slide_layouts[-1]
    new_slide = prs.slides.add_slide(blank_layout)

    # 2) XML层面复制 shape
    src_spTree = src_slide._element.xpath('./p:cSld/p:spTree')[0]
    new_spTree = new_slide._element.xpath('./p:cSld/p:spTree')[0]

    # 先清空 new_slide 里除了背景(bg)以外的元素
    for shape_elem in list(new_spTree):
        if shape_elem.tag.endswith('bg'):
            continue
        new_spTree.remove(shape_elem)

    # 把 src_slide 的 spTree 整体复制过来
    for elem in src_spTree:
        new_spTree.append(copy.deepcopy(elem))

    # 3) 将新幻灯片移动到 src_slide 后面
    src_idx = get_slide_index(prs, src_slide)
    if src_idx == -1:
        # 理论上不该发生, 如果 src_slide 来自同一个 prs
        print("警告: 未能找到 src_slide 的下标, 无法定位移动. 新slide保持末尾.")
        return new_slide
    # 新增幻灯片当前在末尾
    new_idx = len(prs.slides) - 1
    # 我们想让它紧挨在src_slide后面 => 目标索引 = src_idx + 1
    target_idx = src_idx + 1

    # 若 target_idx < new_idx, 说明要往前移动
    if target_idx < new_idx:
        move_slide(prs, new_idx, target_idx)

    return new_slide

def get_slide_index(prs: Presentation, slide: Slide) -> int:
    """
    在 prs.slides 中找到给定 slide 的下标(0-based)。
    如果没找到，返回 -1。
    """
    for i, s in enumerate(prs.slides):
        if s == slide:
            return i
    return -1

def move_slide(prs: Presentation, from_idx: int, to_idx: int):
    """
    将 prs.slides[from_idx] 移动到下标 to_idx(0-based)。
    内部通过操作 _sldIdLst 实现顺序插拔。
    注意：这属于 python-pptx 的私有结构，升级或改动可能影响。
    """
    sldIdLst = prs.slides._sldIdLst  # SlideIdList
    slide_id = sldIdLst[from_idx]
    sldIdLst.remove(slide_id)
    sldIdLst.insert(to_idx, slide_id)